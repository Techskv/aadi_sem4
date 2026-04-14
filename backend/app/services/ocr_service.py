"""
OCR Service - Text extraction from images and PDFs.

Priority for handwritten / scanned content:
  1. Digital PDF  → PyPDF2 (instant, no deps)
  2. Scanned/Handwritten PDF → Vision LLM (Llama) — best accuracy for handwriting
  3. Fallback → Tesseract OCR (if Vision LLM unavailable)

Supported file types: pdf, png, jpg, jpeg, tiff, txt, ppt, pptx
"""
import os
import base64
import logging
import tempfile
from typing import List

from app.config import get_settings

settings = get_settings()
logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────────────────────────────
# Dependency availability checks
# ─────────────────────────────────────────────────────────────────────────────

# PyPDF2 — digital text extraction (no OCR)
PYPDF2_AVAILABLE = False
try:
    from PyPDF2 import PdfReader
    PYPDF2_AVAILABLE = True
except ImportError:
    logger.warning("PyPDF2 not available. Digital PDF extraction disabled.")

# pdf2image + Poppler (Deprecated — PyMuPDF preferred)
PDF2IMAGE_AVAILABLE = False
try:
    from pdf2image import convert_from_path
    # We only mark it available if it's actually usable (check Poppler)
    if getattr(settings, "POPPLER_PATH", ""):
        PDF2IMAGE_AVAILABLE = True
except ImportError:
    pass

# PyMuPDF (fitz) — Better, self-contained PDF to image conversion
PYMUPDF_AVAILABLE = False
try:
    import fitz
    PYMUPDF_AVAILABLE = True
except ImportError:
    logger.warning("PyMuPDF not available. Install it: pip install pymupdf")

# Pillow — image manipulation / base64 encoding
PILLOW_AVAILABLE = False
try:
    from PIL import Image
    import io
    PILLOW_AVAILABLE = True
except ImportError:
    logger.warning("Pillow not available. Install it: pip install Pillow")

# Vision LLM — Vision model for handwriting OCR (Llama via OpenAI-compatible API)
VISION_LLM_AVAILABLE = False
vision_client = None
try:
    from openai import OpenAI
    if settings.LLM_API_KEY and settings.LLM_BASE_URL:
        vision_client = OpenAI(
            api_key=settings.LLM_API_KEY,
            base_url=settings.LLM_BASE_URL,
        )
        VISION_LLM_AVAILABLE = True
        logger.info(f"Vision LLM client initialized (base_url={settings.LLM_BASE_URL})")
    else:
        logger.warning("LLM_API_KEY or LLM_BASE_URL not set — Vision OCR disabled.")
except ImportError:
    logger.warning("openai package not installed — Vision OCR disabled.")
except Exception as e:
    logger.warning(f"Vision LLM init failed: {e}")

# Tesseract — traditional OCR fallback
TESSERACT_AVAILABLE = False
try:
    import pytesseract
    import cv2
    import numpy as np
    tesseract_paths = [
        settings.TESSERACT_CMD,
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        "tesseract",
    ]
    for path in tesseract_paths:
        try:
            pytesseract.pytesseract.tesseract_cmd = path
            pytesseract.get_tesseract_version()
            TESSERACT_AVAILABLE = True
            logger.info(f"Tesseract found at: {path}")
            break
        except Exception:
            continue
    if not TESSERACT_AVAILABLE:
        logger.warning("Tesseract OCR not found. Vision LLM will be used instead.")
except ImportError:
    logger.warning("pytesseract/cv2 not installed.")

# python-pptx — PowerPoint text extraction
PPTX_AVAILABLE = False
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    logger.warning("python-pptx not available.")

# ─────────────────────────────────────────────────────────────────────────────
# Helper: Poppler path for pdf2image on Windows
# ─────────────────────────────────────────────────────────────────────────────

def _get_poppler_path():
    """Return the configured Poppler binary path, or None to use system PATH."""
    p = getattr(settings, "POPPLER_PATH", "").strip()
    return p if p else None


# ─────────────────────────────────────────────────────────────────────────────
# Helper: Convert PIL image to base64 PNG string
# ─────────────────────────────────────────────────────────────────────────────

def _pil_to_base64(pil_image) -> str:
    """Encode a PIL image as a base64 PNG string."""
    buf = io.BytesIO()
    pil_image.save(buf, format="PNG")
    return base64.b64encode(buf.getvalue()).decode("utf-8")


# ─────────────────────────────────────────────────────────────────────────────
# OCR Service
# ─────────────────────────────────────────────────────────────────────────────

class OCRService:
    """
    Service for text extraction from answer sheets.

    Extraction pipeline:
      PDF → try PyPDF2 (digital) → if blank, try Vision LLM → fallback Tesseract
      Image → try Vision LLM → fallback Tesseract
    """

    # ── Vision model (Llama via OpenAI-compatible API) ────────────────────────
    @property
    def VISION_MODEL(self) -> str:
        """Returns the vision model from settings."""
        return getattr(settings, "LLM_VISION_MODEL", "llama-3.2-11b-vision-preview")

    def extract_text_via_vision_llm(self, images: list) -> str:
        """
        Sends images to Llama vision model (via OpenAI-compatible API) for OCR.
        Handles handwritten and scanned content.
        """
        page_texts = []

        for i, pil_img in enumerate(images):
            logger.info(f"Vision OCR: processing page/image {i + 1}/{len(images)}")

            # Resize large images for API limits
            max_dim = 2048
            w, h = pil_img.size
            if max(w, h) > max_dim:
                scale = max_dim / max(w, h)
                pil_img = pil_img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

            b64_img = _pil_to_base64(pil_img)
            text = None

            if VISION_LLM_AVAILABLE and vision_client:
                try:
                    logger.info(f"  Using Vision LLM: {self.VISION_MODEL}")
                    response = vision_client.chat.completions.create(
                        model=self.VISION_MODEL,
                        messages=[
                            {
                                "role": "user",
                                "content": [
                                    {
                                        "type": "image_url",
                                        "image_url": {
                                            "url": f"data:image/png;base64,{b64_img}"
                                        },
                                    },
                                    {
                                        "type": "text",
                                        "text": (
                                            "You are an expert OCR engine. Transcribe the FULL content "
                                            "exactly as written, preserving all structure and numbering. "
                                            "No summaries. Only raw text."
                                        ),
                                    },
                                ],
                            }
                        ],
                        max_tokens=4096,
                        temperature=0,
                    )
                    text = response.choices[0].message.content or ""
                    logger.info(f"  Success via Vision LLM ({len(text)} chars)")
                except Exception as e:
                    logger.error(f"  Vision LLM failed: {e}")

            if text:
                page_texts.append(text.strip())
            else:
                # If vision fails, we want the caller (process_submission) to know
                # so it can fallback to Tesseract if available.
                raise RuntimeError(f"Vision OCR failed on page {i + 1}. Check LLM_API_KEY.")

        return "\n\n--- PAGE BREAK ---\n\n".join(page_texts)

    # ── PDF helpers ───────────────────────────────────────────────────────────

    def extract_text_from_pdf_digital(self, pdf_path: str) -> str:
        """Extract selectable text from a digital PDF using PyPDF2."""
        if not PYPDF2_AVAILABLE:
            raise RuntimeError("PyPDF2 is not installed.")
        reader = PdfReader(pdf_path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append(text.strip() if text.strip() else f"[Page {i + 1}: no digital text]")
        return "\n\n--- PAGE BREAK ---\n\n".join(pages)

    def _pdf_to_pil_images_fitz(self, pdf_path: str) -> list:
        """Convert all PDF pages to PIL images using PyMuPDF (no Poppler needed)."""
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF (fitz) is not installed.")
        
        doc = fitz.open(pdf_path)
        images = []
        for page in doc:
            # 300 DPI (72 * 4.166... = 300)
            pix = page.get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            images.append(img)
        doc.close()
        return images

    def _pdf_to_pil_images(self, pdf_path: str) -> list:
        """Convert all PDF pages to PIL images. Prefers PyMuPDF (fitz)."""
        if PYMUPDF_AVAILABLE:
            try:
                return self._pdf_to_pil_images_fitz(pdf_path)
            except Exception as e:
                logger.warning(f"PyMuPDF rendering failed: {e}")

        if PDF2IMAGE_AVAILABLE:
            try:
                poppler = _get_poppler_path()
                return convert_from_path(pdf_path, dpi=300, poppler_path=poppler)
            except Exception as e:
                logger.warning(f"pdf2image fallback failed: {e}")
        
        raise RuntimeError(
            "No PDF-to-image converter available. "
            "Please install PyMuPDF (pip install pymupdf) for best results."
        )

    def extract_text_from_pdf_vision(self, pdf_path: str) -> str:
        """Convert PDF to images, then OCR via Vision LLM (handles handwriting)."""
        images = self._pdf_to_pil_images(pdf_path)
        return self.extract_text_via_vision_llm(images)

    def extract_text_from_pdf_tesseract(self, pdf_path: str) -> str:
        """Fallback: OCR via Tesseract (printed text only, poor on handwriting)."""
        if not TESSERACT_AVAILABLE:
            raise RuntimeError("Tesseract OCR is not installed.")
        images = self._pdf_to_pil_images(pdf_path)
        results = []
        for i, pil_image in enumerate(images):
            img_np = np.array(pil_image)
            img_bgr = cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR)
            processed = self._preprocess_image(img_bgr)
            import io as _io
            pil_proc = Image.fromarray(processed)
            text = pytesseract.image_to_string(pil_proc, config="--oem 3 --psm 6", lang="eng")
            results.append(text.strip())
        return "\n\n--- PAGE BREAK ---\n\n".join(results)

    # ── Image helpers ─────────────────────────────────────────────────────────

    def extract_text_from_image_vision(self, image_path: str) -> str:
        """OCR a single image via Vision LLM (supports handwriting)."""
        if not PILLOW_AVAILABLE:
            raise RuntimeError("Pillow is required.")
        pil_img = Image.open(image_path).convert("RGB")
        return self.extract_text_via_vision_llm([pil_img])

    def extract_text_from_image_tesseract(self, image_path: str) -> str:
        """OCR a single image via Tesseract (fallback)."""
        if not TESSERACT_AVAILABLE:
            raise RuntimeError("Tesseract OCR is not installed.")
        image = cv2.imread(image_path)
        if image is None:
            raise ValueError(f"Could not read image: {image_path}")
        processed = self._preprocess_image(image)
        pil_img = Image.fromarray(processed)
        return pytesseract.image_to_string(pil_img, config="--oem 3 --psm 6", lang="eng").strip()

    # ── PPTX ─────────────────────────────────────────────────────────────────

    def extract_text_from_pptx(self, pptx_path: str) -> str:
        """Extract text from a PowerPoint file."""
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx is not installed.")
        prs = Presentation(pptx_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            slides.append("\n".join(texts) if texts else f"[Slide {i + 1}: No extractable text]")
        return "\n\n--- SLIDE BREAK ---\n\n".join(slides)

    # ── Plain text ────────────────────────────────────────────────────────────

    def extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from a plain text file."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()

    # ── Main entry point ──────────────────────────────────────────────────────

    def process_submission(self, file_path: str, file_type: str) -> str:
        """
        Main entry point — picks the best available extraction method.

        Pipeline for PDF:
          1. PyPDF2 digital extraction (fast, works for typed PDFs)
          2. If result is blank/stub → Vision LLM (best for handwriting)
          3. If Vision LLM unavailable → Tesseract OCR (printed text only)

        Pipeline for images:
          1. Vision LLM
          2. Tesseract fallback
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_type = file_type.lower().strip(".")

        # ── Plain text ──────────────────────────────────────────────────────
        if file_type == "txt":
            return self.extract_text_from_txt(file_path)

        # ── PowerPoint ─────────────────────────────────────────────────────
        if file_type in ("ppt", "pptx"):
            return self.extract_text_from_pptx(file_path)

        # ── PDF ─────────────────────────────────────────────────────────────
        if file_type == "pdf":
            digital_text = ""

            # Step 1: Try digital extraction
            if PYPDF2_AVAILABLE:
                try:
                    digital_text = self.extract_text_from_pdf_digital(file_path)
                    # Strip ALL placeholder and separator lines before judging real content
                    meaningful = "\n".join(
                        line for line in digital_text.splitlines()
                        if line.strip()
                        and not line.startswith("[Page")
                        and not line.startswith("---")
                    ).strip()
                    if meaningful and len(meaningful) > 30:
                        logger.info(f"PDF: digital text extracted ({len(meaningful)} chars) via PyPDF2")
                        return digital_text
                    else:
                        logger.info(
                            f"PDF: only {len(meaningful)} meaningful chars from PyPDF2 "
                            f"(scanned/handwritten) — switching to Vision LLM OCR"
                        )
                        digital_text = ""   # discard stubs/separators

                except Exception as e:
                    logger.warning(f"PyPDF2 failed: {e}")

            # Step 2: Vision LLM OCR (handwriting support ✅)
            if VISION_LLM_AVAILABLE and (PYMUPDF_AVAILABLE or PDF2IMAGE_AVAILABLE) and PILLOW_AVAILABLE:
                try:
                    logger.info("PDF: using Vision LLM OCR for handwritten/scanned content")
                    return self.extract_text_from_pdf_vision(file_path)
                except Exception as e:
                    logger.warning(f"Vision LLM PDF OCR failed: {e}")

            # Step 3: Tesseract fallback
            if TESSERACT_AVAILABLE and (PYMUPDF_AVAILABLE or PDF2IMAGE_AVAILABLE):
                try:
                    logger.info("PDF: falling back to Tesseract OCR")
                    return self.extract_text_from_pdf_tesseract(file_path)
                except Exception as e:
                    logger.warning(f"Tesseract PDF OCR failed: {e}")

            raise RuntimeError(
                "Cannot extract text from this PDF.\n"
                "• For handwritten PDFs: Vision LLM is required (check LLM_API_KEY + PyMuPDF)\n"
                "• For scanned printed PDFs: install Tesseract + PyMuPDF\n"
                "• For typed PDFs: install PyPDF2"
            )


        # ── Images ─────────────────────────────────────────────────────────
        if file_type in ("png", "jpg", "jpeg", "tiff", "bmp", "webp"):
            # Prefer Vision LLM for handwriting
            if VISION_LLM_AVAILABLE and PILLOW_AVAILABLE:
                try:
                    logger.info(f"Image: using Vision LLM OCR ({file_type})")
                    return self.extract_text_from_image_vision(file_path)
                except Exception as e:
                    logger.warning(f"Vision LLM image OCR failed: {e}")

            # Tesseract fallback
            if TESSERACT_AVAILABLE:
                logger.info(f"Image: using Tesseract OCR fallback ({file_type})")
                return self.extract_text_from_image_tesseract(file_path)

            raise RuntimeError(
                "Cannot extract text from this image.\n"
                "• Install openai package and set LLM_API_KEY for handwriting support.\n"
                "• Or install Tesseract for printed-text images."
            )

        raise ValueError(f"Unsupported file type: {file_type}")

    # ── Image preprocessing (for Tesseract) ──────────────────────────────────

    def _preprocess_image(self, image) -> "np.ndarray":
        """Standard preprocessing to improve Tesseract accuracy on scanned images."""
        try:
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY) if len(image.shape) == 3 else image
            denoised = cv2.bilateralFilter(gray, 9, 75, 75)
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
            enhanced = clahe.apply(denoised)
            binary = cv2.adaptiveThreshold(
                enhanced, 255,
                cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                cv2.THRESH_BINARY,
                11, 2,
            )
            return binary
        except Exception as e:
            logger.error(f"Image preprocessing error: {e}")
            return image


ocr_service = OCRService()
